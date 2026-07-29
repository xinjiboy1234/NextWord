# -*- coding: utf-8 -*-
"""生成 NextWord 产品功能介绍 PPT（给领导汇报用）。

用法（Windows）：
    report/.venv/Scripts/python.exe report/build-ppt.py

版式：16:9（13.33 x 7.5 英寸），白底，黑白主色 + 深青点缀，中文字体微软雅黑。
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 常量
FONT = "微软雅黑"
ACCENT = RGBColor(0x0F, 0x76, 0x6E)        # 深青
ACCENT_LIGHT = RGBColor(0xE6, 0xF4, 0xF2)  # 深青浅色块
DARK = RGBColor(0x1F, 0x1F, 0x1F)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT_GRAY = RGBColor(0xF3, 0xF3, 0xF3)
LINE_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = 1.45   # 内容区起始（英寸）
FOOTER_Y = 7.08

OUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "NextWord-产品功能介绍.pptx")


# ---------------------------------------------------------------- 基础工具
def _style_run(run, size, bold=False, color=DARK):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = FONT
    # 同步设置东亚字体，保证中文用微软雅黑
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
        rPr.append(rPr.makeelement(qn(tag), {"typeface": FONT}))


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15, space_after=6):
    """lines: [(text, size, bold, color), ...] 或 [(text,size,bold,color,align), ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        text, size, bold, color = item[0], item[1], item[2], item[3]
        p_align = item[4] if len(item) > 4 else align
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = p_align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        _style_run(run, size, bold, color)
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_bullets(slide, x, y, w, h, items, size=16, gap=8, marker_color=ACCENT,
                text_color=DARK):
    """items: [str] 或 [(text, bold)]，每段一个要点。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        text, bold = (item, False) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.2
        p.space_after = Pt(gap)
        m = p.add_run()
        m.text = "▪ "
        _style_run(m, size, True, marker_color)
        r = p.add_run()
        r.text = text
        _style_run(r, size, bold, text_color)
    return tb


def add_arrow(slide, x, y, w, h, direction="right", fill=ACCENT):
    shape = {"right": MSO_SHAPE.RIGHT_ARROW,
             "left": MSO_SHAPE.LEFT_ARROW,
             "down": MSO_SHAPE.DOWN_ARROW}[direction]
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def stat_box(slide, x, y, w, h, number, label, note=None):
    add_rect(slide, x, y, w, h, LIGHT_GRAY)
    add_rect(slide, x, y, w, Inches(0.07), ACCENT)
    lines = [(number, 30, True, ACCENT), (label, 13, True, DARK)]
    if note:
        lines.append((note, 11, False, GRAY))
    add_text(slide, x + Inches(0.2), y + Inches(0.22), w - Inches(0.4),
             h - Inches(0.35), lines, line_spacing=1.1, space_after=4)


# ---------------------------------------------------------------- 版式骨架
def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def content_slide(prs, title, page_no, kicker=None):
    """统一内容页版式：标题条 + 分隔线 + 页脚页码。"""
    slide = blank_slide(prs)
    add_rect(slide, MARGIN, Inches(0.42), Inches(0.09), Inches(0.5), ACCENT)
    if kicker:
        add_text(slide, Inches(0.8), Inches(0.28), Inches(11.8), Inches(0.3),
                 [(kicker, 12, True, ACCENT)])
        add_text(slide, Inches(0.8), Inches(0.52), Inches(11.8), Inches(0.5),
                 [(title, 26, True, DARK)])
    else:
        add_text(slide, Inches(0.8), Inches(0.38), Inches(11.8), Inches(0.55),
                 [(title, 28, True, DARK)], anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, MARGIN, Inches(1.18), SLIDE_W - MARGIN * 2, Pt(1.2), LINE_GRAY)
    add_text(slide, MARGIN, Inches(FOOTER_Y), Inches(6), Inches(0.3),
             [("NextWord · 产品功能介绍与迭代成果汇报", 9, False, GRAY)])
    add_text(slide, SLIDE_W - MARGIN - Inches(0.6), Inches(FOOTER_Y),
             Inches(0.6), Inches(0.3),
             [("%d / 14" % page_no, 9, False, GRAY)], align=PP_ALIGN.RIGHT)
    return slide


# ---------------------------------------------------------------- 各页
def slide_cover(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, 0, Inches(2.2), Inches(0.35), Inches(3.1), ACCENT)
    add_rect(slide, Inches(0.35), Inches(2.2), Pt(2), Inches(3.1), LINE_GRAY)
    add_text(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(0.4),
             [("产品功能介绍与迭代成果汇报", 16, True, ACCENT)])
    add_text(slide, Inches(1.0), Inches(2.35), Inches(11.5), Inches(1.8),
             [("NextWord", 54, True, DARK),
              ("AI 驱动的英语表达能力训练", 34, True, DARK)],
             line_spacing=1.15, space_after=10)
    add_text(slide, Inches(1.0), Inches(4.35), Inches(11), Inches(0.8),
             [("从「背单词工具」到「以表达能力为核心」的 AI-native 学习产品",
               16, False, GRAY)])
    add_text(slide, Inches(1.0), Inches(6.5), Inches(11), Inches(0.4),
             [("汇报时间：2026-07-29    ｜    迭代进度：I1 – I6 六轮迭代全部通过验收",
               13, False, GRAY)])


def slide_positioning(prs):
    slide = content_slide(prs, "从「背单词工具」到「表达能力训练」", 2, kicker="01 · 产品定位")
    # 左右对比
    y = Inches(CONTENT_TOP)
    h = Inches(1.9)
    add_rect(slide, MARGIN, y, Inches(5.3), h, LIGHT_GRAY)
    add_text(slide, MARGIN + Inches(0.3), y + Inches(0.25), Inches(4.7), Inches(1.5),
             [("过去 · 词汇工具", 14, True, GRAY),
              ("以「认识单词」为目标", 16, True, DARK),
              ("背词、拼写为主，学到的是识别量", 13, False, GRAY)],
             space_after=8)
    add_arrow(slide, Inches(5.98), y + Inches(0.75), Inches(0.75), Inches(0.45))
    add_rect(slide, Inches(6.85), y, Inches(5.93), h, ACCENT_LIGHT)
    add_text(slide, Inches(7.15), y + Inches(0.25), Inches(5.3), Inches(1.5),
             [("现在 · AI-native 学习产品", 14, True, ACCENT),
              ("以「表达能力」为核心", 16, True, DARK),
              ("测评、画像、计划、洞察全部围绕「能不能用出来」", 13, False, DARK)],
             space_after=8)
    # 两个核心观点
    y2 = Inches(3.7)
    add_rect(slide, MARGIN, y2, Inches(5.95), Inches(2.9), WHITE, line=LINE_GRAY)
    add_rect(slide, MARGIN, y2, Inches(5.95), Inches(0.5), DARK)
    add_text(slide, MARGIN + Inches(0.25), y2 + Inches(0.08), Inches(5.5), Inches(0.35),
             [("核心观点一：等级是外壳，画像是内核", 14, True, WHITE)])
    add_bullets(slide, MARGIN + Inches(0.3), y2 + Inches(0.75), Inches(5.4), Inches(2.0),
                ["固定等级 / 分数只是给用户看的外壳",
                 "真正的内核是详细的画像与评价：WeaknessProfile → 交叉验证 → 定制内容 → 瓶颈洞察 → 重规划",
                 "规则引擎掌握分数权威，AI 不直接改分"], size=14, gap=10)
    add_rect(slide, Inches(6.83), y2, Inches(5.95), Inches(2.9), WHITE, line=LINE_GRAY)
    add_rect(slide, Inches(6.83), y2, Inches(5.95), Inches(0.5), ACCENT)
    add_text(slide, Inches(7.08), y2 + Inches(0.08), Inches(5.5), Inches(0.35),
             [("核心观点二：词的毕业标准是「能用」", 14, True, WHITE)])
    add_bullets(slide, Inches(7.13), y2 + Inches(0.75), Inches(5.4), Inches(2.0),
                ["词的掌握不以「认识」为终点",
                 "四阶段毕业路径：认识 → 回忆 → 造句使用 → 自发使用",
                 "产出任务只用水平带内的词，阅读靠查词机制兜底"], size=14, gap=10)


def slide_loop(prs):
    slide = content_slide(prs, "核心学习闭环", 3, kicker="02 · 产品概览")
    boxes = ["测评定级", "弱项画像", "7 日定制计划", "日常学习", "瓶颈洞察", "自动重规划"]
    subs = ["产出型自适应测评", "WeaknessProfile\n+ Verifier 核查",
            "Planner 夜间生成", "背词 / 拼写\n造句 / 阅读",
            "规则筛查 +\nInsightAgent", "性质变化即触发\n每周兜底"]
    bw, bh, aw = Inches(3.35), Inches(1.55), Inches(0.85)
    x0 = MARGIN
    y1, y2 = Inches(1.7), Inches(4.35)
    xs = [x0, x0 + bw + aw, x0 + (bw + aw) * 2]
    # 第一行：0 → 1 → 2
    for i in range(3):
        _flow_box(slide, xs[i], y1, bw, bh, boxes[i], subs[i], i + 1)
        if i < 2:
            add_arrow(slide, xs[i] + bw + Inches(0.12), y1 + Inches(0.55),
                      aw - Inches(0.24), Inches(0.45))
    # 右侧下行箭头
    add_arrow(slide, xs[2] + bw / 2 - Inches(0.22), y1 + bh + Inches(0.12),
              Inches(0.45), y2 - y1 - bh - Inches(0.24), direction="down")
    # 第二行：从右向左 3 → 4 → 5（日常学习在右，自动重规划在左）
    positions = {3: xs[2], 4: xs[1], 5: xs[0]}
    for idx in (3, 4, 5):
        _flow_box(slide, positions[idx], y2, bw, bh, boxes[idx], subs[idx], idx + 1)
    add_arrow(slide, xs[2] - aw + Inches(0.12), y2 + Inches(0.55),
              aw - Inches(0.24), Inches(0.45), direction="left")
    add_arrow(slide, xs[1] - aw + Inches(0.12), y2 + Inches(0.55),
              aw - Inches(0.24), Inches(0.45), direction="left")
    # 左端回环：自动重规划 → 回到计划/学习
    add_rect(slide, xs[0] + Inches(0.35), y1 + bh + Inches(0.3),
             Pt(2.5), y2 - y1 - bh - Inches(0.6), ACCENT)
    add_text(slide, xs[0] + Inches(0.5), Inches(3.35), Inches(4.6), Inches(0.7),
             [("↺ 瓶颈性质变化 → 重生成画像 → 新的 7 日计划", 12, True, ACCENT)])
    add_text(slide, MARGIN, Inches(6.25), Inches(12.2), Inches(0.6),
             [("闭环要点：每一步都由已验证的画像驱动；规则引擎保证分数确定，AI 负责解读与规划——用户在「计划 → 执行 → 洞察 → 重规划」中持续进阶。",
               13, False, GRAY)])


def _flow_box(slide, x, y, w, h, title, sub, num):
    add_rect(slide, x, y, w, h, LIGHT_GRAY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(slide, x, y, Inches(0.09), h, ACCENT)
    sub_lines = [(t, 11, False, GRAY) for t in sub.split("\n")]
    add_text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.45), h - Inches(0.3),
             [("%d. %s" % (num, title), 16, True, DARK)] + sub_lines,
             line_spacing=1.1, space_after=3)


def slide_i1(prs):
    slide = content_slide(prs, "场景内容建设：先解决「学什么」", 4, kicker="03 · 迭代一（I1）")
    y = Inches(CONTENT_TOP)
    bw, bh, gap = Inches(2.9), Inches(1.7), Inches(0.21)
    stat_box(slide, MARGIN, y, bw, bh, "7 × 20", "生活表达场景", "两层 taxonomy，7 大类覆盖 20 个子场景")
    stat_box(slide, MARGIN + (bw + gap), y, bw, bh, "1520", "内置场景词库", "LLM 生成 + 人工规则验收口径守护")
    stat_box(slide, MARGIN + (bw + gap) * 2, y, bw, bh, "逐词标注", "scenario / utility / role", "每个词带场景、实用度与角色标签")
    stat_box(slide, MARGIN + (bw + gap) * 3, y, bw, bh, "98%", "抽样标注准确率", "按子场景抽样人工核对")
    y2 = Inches(3.55)
    add_rect(slide, MARGIN, y2, Inches(5.95), Inches(3.0), LIGHT_GRAY)
    add_text(slide, MARGIN + Inches(0.3), y2 + Inches(0.2), Inches(5.4), Inches(0.4),
             [("场景覆盖（7 大类）", 15, True, ACCENT)])
    add_bullets(slide, MARGIN + Inches(0.3), y2 + Inches(0.7), Inches(5.4), Inches(2.2),
                ["日常起居、出行交通、购物点餐、社交沟通",
                 "工作职场、学习校园、健康医疗等真实生活语境",
                 "情境化测评与造句题的场景素材均取自该 taxonomy"], size=14, gap=10)
    add_rect(slide, Inches(6.83), y2, Inches(5.95), Inches(3.0), LIGHT_GRAY)
    add_text(slide, Inches(7.13), y2 + Inches(0.2), Inches(5.4), Inches(0.4),
             [("为什么重要", 15, True, ACCENT)])
    add_bullets(slide, Inches(7.13), y2 + Inches(0.7), Inches(5.4), Inches(2.2),
                ["词不再是孤立列表，而是挂在「能用上的场景」下",
                 "utility=low 的词直接不入库，保证内容质量",
                 "为 I2 情境测评、I3 主攻场景计划提供内容地基"], size=14, gap=10)


def slide_i2_assessment(prs):
    slide = content_slide(prs, "测评重构：产出型为主的自适应测评", 5, kicker="04 · 迭代二（I2）")
    y = Inches(CONTENT_TOP)
    # 左侧大块：60% 产出题
    add_rect(slide, MARGIN, y, Inches(3.7), Inches(4.9), ACCENT)
    add_text(slide, MARGIN + Inches(0.35), y + Inches(0.5), Inches(3.0), Inches(3.6),
             [("60%", 60, True, WHITE),
              ("产出题占比", 18, True, WHITE),
              ("造句 + 情境表达为主", 13, False, WHITE),
              ("识别题仅作参考展示", 13, False, WHITE),
              ("不参与定级", 13, False, WHITE)], space_after=10)
    # 右侧要点
    x2 = Inches(4.6)
    add_bullets(slide, x2, y + Inches(0.1), Inches(8.1), Inches(2.2),
                [("产出题全部走 LLM 真实四维评分：语法 / 自然度 / 词汇 / 相关度，复用造句工作室评分链路，词数启发式已废弃", False),
                 ("自适应分块收敛：每块 5 题，块表现决定升带 / 降带，2–3 块收敛，总题量 ≤ 15", False),
                 ("主定级 = 表达力综合分（四维加权），与 Score 分带对齐、封顶 C1", False)], size=15, gap=12)
    y3 = Inches(4.15)
    add_rect(slide, x2, y3, Inches(8.1), Inches(2.2), LIGHT_GRAY)
    add_text(slide, x2 + Inches(0.3), y3 + Inches(0.2), Inches(7.5), Inches(0.4),
             [("出题纪律（词池管控）", 14, True, ACCENT)])
    add_bullets(slide, x2 + Inches(0.3), y3 + Inches(0.65), Inches(7.5), Inches(1.4),
                ["出题词只选水平带内且 utility=high / medium，绝不超带",
                 "情境场景取自 I1 taxonomy；阅读题从库内分级短文就近选文，答案位置随机"], size=13, gap=8)


def slide_i2_profile(prs):
    slide = content_slide(prs, "WeaknessProfile + Verifier：让 AI 结论可信", 6, kicker="04 · 迭代二（I2）")
    y = Inches(CONTENT_TOP)
    add_rect(slide, MARGIN, y, Inches(5.95), Inches(3.5), WHITE, line=LINE_GRAY)
    add_rect(slide, MARGIN, y, Inches(5.95), Inches(0.5), DARK)
    add_text(slide, MARGIN + Inches(0.25), y + Inches(0.08), Inches(5.5), Inches(0.35),
             [("WeaknessProfile：每条 Finding 五要素", 14, True, WHITE)])
    add_bullets(slide, MARGIN + Inches(0.3), y + Inches(0.7), Inches(5.4), Inches(2.6),
                ["维度：scenario / skill / reading",
                 "强弱：strength / weakness / neutral",
                 "结论文案：一句可读的结论",
                 "证据引用：指向库内真实学习记录",
                 "置信度：high / medium / low"], size=14, gap=9)
    x2 = Inches(6.83)
    add_rect(slide, x2, y, Inches(5.95), Inches(3.5), WHITE, line=LINE_GRAY)
    add_rect(slide, x2, y, Inches(5.95), Inches(0.5), ACCENT)
    add_text(slide, x2 + Inches(0.25), y + Inches(0.08), Inches(5.5), Inches(0.35),
             [("Verifier：不调 LLM 的逐条机械核查", 14, True, WHITE)])
    add_bullets(slide, x2 + Inches(0.3), y + Inches(0.7), Inches(5.4), Inches(2.6),
                ["证据真实存在，且属于本人（防编造、防越权）",
                 "引用数值与库内重算值一致（可重算）",
                 "证据条数支撑置信度（high≥3 / medium≥2 / low≥1）",
                 "任一不通过 → 标记 Questioned 并留原因"], size=14, gap=9)
    y2 = Inches(5.3)
    add_rect(slide, MARGIN, y2, Inches(12.23), Inches(1.25), ACCENT_LIGHT)
    add_text(slide, MARGIN + Inches(0.35), y2 + Inches(0.2), Inches(11.6), Inches(0.9),
             [("关键设计：存疑条目「不展示、不进规划」。", 17, True, ACCENT),
              ("用户看到的每一条画像结论、Planner 消费的每一条输入，都必须先通过 Verifier——这是 AI 结论可信的护城河。",
               13, False, DARK)], space_after=6)


def slide_i3_planner(prs):
    slide = content_slide(prs, "个性化学习计划：每天打开就是「今日任务」", 7, kicker="05 · 迭代三（I3）")
    y = Inches(CONTENT_TOP)
    add_bullets(slide, MARGIN, y, Inches(7.2), Inches(3.4),
                [("Planner 依据已验证的画像（Verified Finding），夜间后台生成 7 日计划", False),
                 ("计划内容：主攻场景（1–2 个）+ 每日词队列 + 阅读推荐（3 篇）+ 每日造句目标（3 词）", False),
                 ("背词、造句、阅读三个模块全部切换为「执行今日计划」，各入口统一 fromPlan 标记", False),
                 ("无计划 / 计划过期自动回退保底策略，用户永远有内容可学", False),
                 ("测评完成自动入队，任务幂等防重复生成", False)], size=15, gap=13)
    x2 = Inches(8.1)
    add_rect(slide, x2, y, Inches(4.68), Inches(4.9), LIGHT_GRAY)
    add_text(slide, x2 + Inches(0.3), y + Inches(0.25), Inches(4.1), Inches(0.4),
             [("接触词机制", 15, True, ACCENT)])
    add_bullets(slide, x2 + Inches(0.3), y + Inches(0.75), Inches(4.1), Inches(3.0),
                ["每日词队列允许 ≤20% 超出水平带的「接触词」",
                 "接触词只要求「认识」，不进产出任务",
                 "在保底之上给用户提供适度挑战与新鲜感"], size=13, gap=10)
    add_rect(slide, x2 + Inches(0.3), y + Inches(3.55), Inches(4.05), Inches(1.0), ACCENT)
    add_text(slide, x2 + Inches(0.5), y + Inches(3.72), Inches(3.7), Inches(0.7),
             [("难度策略分层：背词可掺接触词，产出只用带内词", 12, True, WHITE)],
             line_spacing=1.15)


def slide_i3_insight(prs):
    slide = content_slide(prs, "瓶颈洞察与自动重规划：三层机制", 8, kicker="05 · 迭代三（I3）")
    y = Inches(CONTENT_TOP)
    bw, bh = Inches(3.85), Inches(3.5)
    gap = Inches(0.55)
    layers = [
        ("① 规则筛查", "零 LLM 成本 · 日级运行",
         ["平台期检测", "回避模式识别", "安全词策略识别", "未触发 → 全程零 LLM 调用"]),
        ("② InsightAgent 判定", "事件驱动 · 细读产出原文",
         ["读取近 20 条用户产出原文", "判定 7 类瓶颈性质", "词汇不足 / 组织不成句 / 语法错误多", "表达单调 / 回避 / 中式搭配 / 安全词"]),
        ("③ 自动重规划", "性质变化即触发 + 每周兜底",
         ["瓶颈性质变化 → 重生成画像", "→ 入队 force Planner 重建计划", "性质未变 → 仅记录不打扰", "WeeklyReplanWorker 每周兜底"]),
    ]
    x = MARGIN
    for i, (title, tag, items) in enumerate(layers):
        add_rect(slide, x, y, bw, bh, LIGHT_GRAY)
        add_rect(slide, x, y, bw, Inches(0.07), ACCENT)
        add_text(slide, x + Inches(0.25), y + Inches(0.2), bw - Inches(0.5), Inches(0.75),
                 [(title, 16, True, DARK), (tag, 11, True, ACCENT)], space_after=3)
        add_bullets(slide, x + Inches(0.25), y + Inches(1.05), bw - Inches(0.5), Inches(2.3),
                    items, size=12, gap=6)
        if i < 2:
            add_arrow(slide, x + bw + Inches(0.06), y + Inches(1.5),
                      gap - Inches(0.12), Inches(0.45))
        x += bw + gap
    y2 = Inches(5.45)
    add_rect(slide, MARGIN, y2, Inches(12.23), Inches(1.1), ACCENT_LIGHT)
    add_text(slide, MARGIN + Inches(0.35), y2 + Inches(0.18), Inches(11.6), Inches(0.8),
             [("洞察只影响解读与规划，不改任何分数。", 15, True, ACCENT),
              ("证据沿用画像纪律：编造 / 越权的证据引用在持久化前机械过滤；同日幂等，成本可控。",
               12, False, DARK)], space_after=5)


def slide_i4(prs):
    slide = content_slide(prs, "词毕业生命周期：掌握度 = 能不能用", 9, kicker="06 · 迭代四（I4）")
    y = Inches(1.75)
    bw, bh, aw = Inches(2.6), Inches(2.1), Inches(0.55)
    stages = [("认识", "25", "看词知义\nrecognition 模式"),
              ("回忆", "50", "看义想词\nrecall 拼写正确"),
              ("造句使用", "75", "指定词造句 A/B 档\n确认 prompted use"),
              ("自发使用", "100", "自由表达中自发用出\n= 真正毕业")]
    x = MARGIN
    for i, (name, score, desc) in enumerate(stages):
        fill = ACCENT if i == 3 else LIGHT_GRAY
        tcolor = WHITE if i == 3 else DARK
        scolor = WHITE if i == 3 else ACCENT
        add_rect(slide, x, y, bw, bh, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        lines = [(name, 17, True, tcolor), (score + " 分", 20, True, scolor)]
        lines += [(t, 11, False, GRAY if i < 3 else WHITE) for t in desc.split("\n")]
        add_text(slide, x + Inches(0.22), y + Inches(0.25), bw - Inches(0.4), bh - Inches(0.4),
                 lines, line_spacing=1.12, space_after=4)
        if i < 3:
            add_arrow(slide, x + bw + Inches(0.06), y + Inches(0.85),
                      aw - Inches(0.12), Inches(0.4))
        x += bw + aw
    add_text(slide, MARGIN, y + bh + Inches(0.15), Inches(12.2), Inches(0.35),
             [("四阶段状态机，掌握度由阶段派生（25 / 50 / 75 / 100）", 13, True, ACCENT)],
             align=PP_ALIGN.CENTER)
    y2 = Inches(4.75)
    add_bullets(slide, MARGIN, y2, Inches(12.2), Inches(1.8),
                [("用户自评（记得 / 忘记）只调整 SM-2 复习排程，不再直接加减掌握度——分数权威收归规则引擎", False),
                 ("阶段推进与回退有明确规则：造句含目标词且评分 A/B → 确认使用；D 档 → 退回回忆阶段重进 SM-2 调度", False),
                 ("毕业留痕可追溯（GraduatedFreeExpressionLogId）；Planner 优先把「待确认」词编入造句目标，推动词走向毕业", False)],
                size=14, gap=11)


# ---------------------------------------------------------------- I5 / I6 素材图
_REPORT_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_TIMELINE = os.path.join(_REPORT_DIR, "..", "demo", "agent-story",
                            "output", "timeline-preview.png")
IMG_PLAN_CARD = os.path.join(_REPORT_DIR, "screenshots", "dashboard-plan-card.png")
IMG_INSIGHT_CARD = os.path.join(_REPORT_DIR, "screenshots",
                                "dashboard-insight-card.png")


def slide_i5(prs):
    slide = content_slide(prs, "Agent 协作演示《林晓的七天》：真实链路全留痕", 10,
                          kicker="07 · 迭代五（I5）")
    y = Inches(CONTENT_TOP)
    bw, bh, gap = Inches(2.9), Inches(1.7), Inches(0.21)
    stat_box(slide, MARGIN, y, bw, bh, "零改动", "真实 API + 真实 LLM", "不改代码不改数据，独立库 + qwen-plus")
    stat_box(slide, MARGIN + (bw + gap), y, bw, bh, "25 次", "LLM 调用全留痕", "记录代理逐条捕获，共 22.8k tokens")
    stat_box(slide, MARGIN + (bw + gap) * 2, y, bw, bh, "4 / 4", "画像 Finding 全 Verified", "Verifier 机械核查逐条通过")
    stat_box(slide, MARGIN + (bw + gap) * 3, y, bw, bh, "69 事件", "交互时间轴", "timeline.html 剧情节点 ↔ LLM 对话双向跳转")
    y2 = Inches(3.45)
    add_rect(slide, MARGIN, y2, Inches(5.95), Inches(3.25), LIGHT_GRAY)
    add_text(slide, MARGIN + Inches(0.3), y2 + Inches(0.2), Inches(5.4), Inches(0.4),
             [("七天故事线：Agent 各就各位", 15, True, ACCENT)])
    add_bullets(slide, MARGIN + Inches(0.3), y2 + Inches(0.68), Inches(5.4), Inches(2.5),
                ["Day 1：测评 2 块收敛定 B2 → 画像 4 Finding 全 Verified → Planner 首计划（新用户覆盖率兜底）",
                 "Day 2–6：正常学习转入回避（连接词率真实构造 1.7 → 0），规则引擎零 LLM 捕获信号",
                 "Day 7：Insight 细读原文独立定性「词汇不足」——不给信号盖章，定性权在 Agent",
                 "性质变化自动重规划：新计划精确消费 Verified 场景 Finding，兜底 → 个性化演进完整呈现"],
                size=12, gap=8)
    # 右侧：timeline.html 预览图（1600×900）
    img_w = Inches(5.8)
    img_h = Inches(5.8 * 900 / 1600)
    slide.shapes.add_picture(IMG_TIMELINE, Inches(6.83), y2, width=img_w, height=img_h)
    add_text(slide, Inches(6.83), y2 + img_h + Inches(0.04), Inches(5.95), Inches(0.3),
             [("▲ demo/agent-story/timeline.html 交互时间轴预览", 11, False, GRAY)],
             align=PP_ALIGN.CENTER)


def slide_i6(prs):
    slide = content_slide(prs, "Agent 价值用户可见：今日计划卡 + 学习洞察卡", 11,
                          kicker="08 · 迭代六（I6）")
    y = Inches(CONTENT_TOP)
    col_w = Inches(5.95)
    # 左列：今日学习计划卡
    add_rect(slide, MARGIN, y, col_w, Inches(0.5), ACCENT)
    add_text(slide, MARGIN + Inches(0.25), y + Inches(0.08), Inches(5.5), Inches(0.35),
             [("今日学习计划卡（T-018）", 14, True, WHITE)])
    add_bullets(slide, MARGIN + Inches(0.15), y + Inches(0.65), Inches(5.7), Inches(1.6),
                ["主攻场景中文名 + 计划进度（第 x / 7 天）",
                 "今日带内词 + 接触词数量、造句目标词",
                 "来源徽章：个性化（Verified 画像驱动）/ 探索期（覆盖率兜底）",
                 "无计划时展示引导文案，加载失败静默降级"], size=12, gap=5)
    img_h = Inches(5.95 * 203 / 528)
    slide.shapes.add_picture(IMG_PLAN_CARD, MARGIN, y + Inches(2.35),
                             width=col_w, height=img_h)
    # 右列：学习洞察卡
    x2 = Inches(6.83)
    add_rect(slide, x2, y, col_w, Inches(0.5), DARK)
    add_text(slide, x2 + Inches(0.25), y + Inches(0.08), Inches(5.5), Inches(0.35),
             [("学习洞察卡（T-019）", 14, True, WHITE)])
    add_bullets(slide, x2 + Inches(0.15), y + Inches(0.65), Inches(5.7), Inches(1.6),
                ["7 类瓶颈性质中文名 + 人话解释",
                 "Agent 结论原文呈现，附触发时间",
                 "已重规划显示「已为你调整学习计划」徽章",
                 "措辞面向用户，内部 id 不进 DOM"], size=12, gap=5)
    slide.shapes.add_picture(IMG_INSIGHT_CARD, x2, y + Inches(2.35),
                             width=col_w, height=img_h)
    y3 = y + Inches(2.35) + img_h + Inches(0.25)
    add_rect(slide, MARGIN, y3, Inches(12.23), Inches(0.68), ACCENT_LIGHT)
    add_text(slide, MARGIN + Inches(0.35), y3 + Inches(0.13), Inches(11.6), Inches(0.55),
             [("后端零改动：两张卡全部消费既有只读端点（planner/current、insights/bottleneck/latest）——Agent 在后台干的活，用户终于看得见。",
               13, True, ACCENT)])


def slide_architecture(prs):
    slide = content_slide(prs, "AI 架构：规则引擎掌分，AI 做智能", 12, kicker="09 · 架构设计")
    y = Inches(CONTENT_TOP)
    bw, bh = Inches(3.85), Inches(4.55)
    gap = Inches(0.34)
    cols = [
        ("规则引擎 · 分数权威", DARK, WHITE, WHITE, WHITE,
         ["SM-2 间隔重复", "Score 内核（0–100 三维）", "等级升级规则", "生命周期掌握度派生",
          "全部确定性代码", "AI 不直接改分"]),
        ("AI 负责智能体验", ACCENT, WHITE, WHITE, WHITE,
         ["单词难度 / 场景标注", "造句与情境表达评分", "WeaknessProfile 解读", "瓶颈性质判定",
          "7 日计划生成", "统一 ILLMProvider 抽象"]),
        ("成本可控 · 可插拔", LIGHT_GRAY, DARK, ACCENT, DARK,
         ["规划类任务日级运行", "任务幂等防重复扣费", "规则筛查零 LLM", "查词多级缓存",
          "Mock 降级零依赖", "OpenAI 兼容端点 + 用户级 BYOK"]),
    ]
    x = MARGIN
    for title, fill, tcolor, scolor, bcolor, items in cols:
        add_rect(slide, x, y, bw, bh, fill)
        add_text(slide, x + Inches(0.25), y + Inches(0.25), bw - Inches(0.5), Inches(0.45),
                 [(title, 16, True, tcolor)])
        add_bullets(slide, x + Inches(0.25), y + Inches(0.85), bw - Inches(0.5), Inches(3.5),
                    items, size=13, gap=9, marker_color=scolor, text_color=bcolor)
        x += bw + gap
    add_text(slide, MARGIN, Inches(6.35), Inches(12.2), Inches(0.55),
             [("设计主线：确定性的事情交给规则引擎，需要判断的事情交给 AI，中间用 Verifier 与幂等任务隔开——既智能、又可信、还省钱。",
               13, False, GRAY)])


def slide_engineering(prs):
    slide = content_slide(prs, "工程与质量：每轮迭代都可验证", 13, kicker="10 · 工程质量")
    y = Inches(CONTENT_TOP)
    bw, bh, gap = Inches(2.9), Inches(1.7), Inches(0.21)
    stat_box(slide, MARGIN, y, bw, bh, "149", "自动化测试全绿", "143 单元 + 6 集成（真实 PostgreSQL）")
    stat_box(slide, MARGIN + (bw + gap), y, bw, bh, "6 / 6", "迭代验收通过", "I1–I6 每轮验收全过")
    stat_box(slide, MARGIN + (bw + gap) * 2, y, bw, bh, "真实 LLM", "端到端实测验收", "每轮走 qwen-plus 真实链路验证")
    stat_box(slide, MARGIN + (bw + gap) * 3, y, bw, bh, "3 角色", "迭代小组工作循环", "产品 / 开发 / 测试闭环协作")
    y2 = Inches(3.55)
    add_rect(slide, MARGIN, y2, Inches(5.95), Inches(3.0), LIGHT_GRAY)
    add_text(slide, MARGIN + Inches(0.3), y2 + Inches(0.2), Inches(5.4), Inches(0.4),
             [("技术栈", 15, True, ACCENT)])
    add_bullets(slide, MARGIN + Inches(0.3), y2 + Inches(0.7), Inches(5.4), Inches(2.2),
                ["后端：.NET 10 + ASP.NET Core Minimal API",
                 "前端：React 19 + TypeScript + Vite + Tailwind 4",
                 "数据库：PostgreSQL 16（EF Core 10 + Npgsql）",
                 "LLM：Microsoft.Extensions.AI，可切 OpenAI 兼容端点"], size=13, gap=9)
    add_rect(slide, Inches(6.83), y2, Inches(5.95), Inches(3.0), LIGHT_GRAY)
    add_text(slide, Inches(7.13), y2 + Inches(0.2), Inches(5.4), Inches(0.4),
             [("工作循环", 15, True, ACCENT)])
    add_bullets(slide, Inches(7.13), y2 + Inches(0.7), Inches(5.4), Inches(2.2),
                ["产品对齐愿景提需求 → 开发实现 → 测试细心验收",
                 "每轮不足记录进任务表，作为下一轮迭代输入",
                 "一个任务一次提交，测试通过 + 构建通过方可交付",
                 "数据库迁移即时生成，schema 演进可追溯",
                 "演示录屏实测发现并修复挑战页选项组件 bug（T-016），一处修复全站同愈"], size=13, gap=8)


def slide_next(prs):
    slide = content_slide(prs, "下一步规划", 14, kicker="11 · Next Steps")
    y = Inches(CONTENT_TOP)
    items = [
        ("阅读助手 Agent 接入前端", "后端能力已就绪，前端接入后阅读场景也有 AI 陪伴"),
        ("E2E 测试纳入 CI", "Playwright 端到端用例进流水线，守住核心学习闭环"),
        ("OpenTelemetry 可观测性", "Agent 调用与后台任务全链路可观测，成本与质量可度量"),
        ("Release Blockers sign-off", "发布前阻断项清单化核对，达到 v1 正式上线标准"),
    ]
    bw, bh = Inches(5.95), Inches(2.15)
    positions = [(MARGIN, y), (Inches(6.83), y),
                 (MARGIN, y + bh + Inches(0.4)), (Inches(6.83), y + bh + Inches(0.4))]
    for i, ((title, desc), (x, yy)) in enumerate(zip(items, positions)):
        add_rect(slide, x, yy, bw, bh, LIGHT_GRAY)
        add_rect(slide, x, yy, Inches(0.09), bh, ACCENT)
        add_text(slide, x + Inches(0.35), yy + Inches(0.3), bw - Inches(0.6), bh - Inches(0.5),
                 [("P%d · %s" % (i, title), 17, True, DARK),
                  (desc, 13, False, GRAY)], line_spacing=1.2, space_after=8)
    add_text(slide, MARGIN, Inches(6.5), Inches(12.2), Inches(0.5),
             [("方向不变：一切迭代继续以「表达能力」为核心，以 VISION-expression-first 为准绳。",
               13, True, ACCENT)], align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- 主流程
def build():
    prs = new_deck()
    slide_cover(prs)
    slide_positioning(prs)
    slide_loop(prs)
    slide_i1(prs)
    slide_i2_assessment(prs)
    slide_i2_profile(prs)
    slide_i3_planner(prs)
    slide_i3_insight(prs)
    slide_i4(prs)
    slide_i5(prs)
    slide_i6(prs)
    slide_architecture(prs)
    slide_engineering(prs)
    slide_next(prs)
    prs.save(OUT_PATH)
    print("saved:", OUT_PATH)


if __name__ == "__main__":
    build()
