# -*- coding: utf-8 -*-
"""验收校验：页数 = 14；形状不超出页面边界；打印每页标题。"""
import os
import sys

from pptx import Presentation
from pptx.util import Emu

sys.stdout.reconfigure(encoding="utf-8")

path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                    "NextWord-产品功能介绍.pptx")
prs = Presentation(path)
W, H = prs.slide_width, prs.slide_height

assert len(prs.slides) == 14, "页数应为 14，实际 %d" % len(prs.slides)
print("页数检查通过：14 页")

problems = []
titles = []
for idx, slide in enumerate(prs.slides, 1):
    best_size, title = -1, None
    for shp in slide.shapes:
        # 边界检查（允许 0.02 英寸容差）
        tol = Emu(int(0.02 * 914400))
        if shp.left < -tol or shp.top < -tol \
                or shp.left + shp.width > W + tol \
                or shp.top + shp.height > H + tol:
            problems.append("第%d页 形状越界: %s left=%.2f top=%.2f right=%.2f bottom=%.2f" % (
                idx, shp.shape_type,
                shp.left / 914400, shp.top / 914400,
                (shp.left + shp.width) / 914400, (shp.top + shp.height) / 914400))
        # 标题 = 标题条区域（内容页 y<1.2in，封面 y<3.0in）内字号最大的文本
        region = Emu(int(3.0 * 914400)) if idx == 1 else Emu(int(1.2 * 914400))
        if shp.has_text_frame and shp.top < region:
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() and run.font.size is not None \
                            and run.font.size > best_size:
                        best_size = run.font.size
                        title = run.text.strip()
    titles.append(title or "(未识别)")

for i, t in enumerate(titles, 1):
    print("第%2d页: %s" % (i, t))

if problems:
    print("\n边界问题：")
    for p in problems:
        print(" -", p)
    sys.exit(1)
print("\n边界检查通过：所有形状均在页面范围内")
print("文件可正常读取：", path)
